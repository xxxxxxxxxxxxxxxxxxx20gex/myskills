import argparse
import shutil
import subprocess
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


EXPORT_CSS = """
html,
body {
  width: 1600px !important;
  height: 900px !important;
  min-height: 900px !important;
  margin: 0 !important;
  padding: 0 !important;
  display: block !important;
  overflow: hidden !important;
}

.slide {
  width: 1600px !important;
  height: 900px !important;
  transform: none !important;
  transform-origin: top left !important;
  box-shadow: none !important;
}
"""


def find_chrome():
    candidates = [
        Path(r"C:\Program Files\Google\Chrome\Application\chrome.exe"),
        Path(r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"),
        Path(r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"),
        Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"),
    ]
    for path in candidates:
        if path.exists():
            return str(path)
    found = shutil.which("chrome") or shutil.which("msedge")
    if found:
        return found
    raise SystemExit("Chrome or Edge was not found.")


def prepare_export_copy(project, temp_root):
    export_project = temp_root / project.name
    ignore = shutil.ignore_patterns(
        "_export_hd_png",
        "_preview*",
        "*.pptx",
        "__pycache__",
    )
    shutil.copytree(project, export_project, ignore=ignore)
    (export_project / "export.css").write_text(EXPORT_CSS, encoding="utf-8")

    slides = sorted(export_project.glob("slide-*.html"))
    if not slides:
        raise SystemExit(f"No slide-*.html files found in {project}")

    for html in slides:
        text = html.read_text(encoding="utf-8")
        marker = "</head>"
        if marker not in text:
            raise SystemExit(f"{html.name} has no </head> marker.")
        text = text.replace(marker, '  <link rel="stylesheet" href="export.css">\n</head>', 1)
        html.write_text(text, encoding="utf-8")
    return export_project


def render_slides(chrome, export_project, preview_dir, scale):
    preview_dir.mkdir(parents=True, exist_ok=True)
    for old in preview_dir.glob("slide-*.png"):
        old.unlink()

    expected = (1600 * scale, 900 * scale)
    rendered = []
    for html in sorted(export_project.glob("slide-*.html")):
        png = preview_dir / f"{html.stem}.png"
        cmd = [
            chrome,
            "--headless=new",
            "--disable-gpu",
            "--hide-scrollbars",
            f"--force-device-scale-factor={scale}",
            "--window-size=1600,900",
            f"--screenshot={png}",
            html.resolve().as_uri(),
        ]
        result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        if result.returncode != 0:
            message = result.stderr.decode("utf-8", "replace")
            raise SystemExit(f"Failed to render {html.name}\n{message}")
        with Image.open(png) as image:
            if image.size != expected:
                raise SystemExit(f"{png.name} rendered at {image.size}, expected {expected}.")
        rendered.append(png)
    return rendered


def build_pptx(images, output):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank = prs.slide_layouts[6]

    for image in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(image), 0, 0, width=prs.slide_width, height=prs.slide_height)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)


def main():
    parser = argparse.ArgumentParser(description="Export fixed-canvas HTML slides to an HD image-based PPTX.")
    parser.add_argument("--project", required=True, help="Folder containing slide-*.html files.")
    parser.add_argument("--output", required=True, help="Output .pptx path.")
    parser.add_argument("--scale", type=int, default=2, choices=[1, 2, 3], help="Chrome device scale factor.")
    parser.add_argument("--preview-dir", help="PNG output folder. Defaults to <project>/_export_hd_png.")
    args = parser.parse_args()

    project = Path(args.project).resolve()
    output = Path(args.output).resolve()
    preview_dir = Path(args.preview_dir).resolve() if args.preview_dir else project / "_export_hd_png"

    chrome = find_chrome()
    with tempfile.TemporaryDirectory(prefix="html_ppt_export_") as temp_dir:
        export_project = prepare_export_copy(project, Path(temp_dir))
        images = render_slides(chrome, export_project, preview_dir, args.scale)
    build_pptx(images, output)

    print(output)
    print(f"slides={len(images)}")
    print(f"bytes={output.stat().st_size}")


if __name__ == "__main__":
    main()
